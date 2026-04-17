from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Amazon brand colors
AMAZON_ORANGE = RGBColor(0xFF, 0x99, 0x00)
AMAZON_NAVY   = RGBColor(0x23, 0x2F, 0x3E)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY     = RGBColor(0x4A, 0x4A, 0x4A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_text(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    set_text(tf, text, size, bold=bold, color=color, align=align, italic=italic)
    return txBox


def add_bullet_box(slide, left, top, width, height, items, size=18,
                   color=AMAZON_NAVY, bold_first=False):
    from pptx.oxml.ns import qn
    import copy
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (icon, label, detail) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)

        run = p.add_run()
        run.text = f"{icon}  {label}"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = AMAZON_NAVY

        if detail:
            p2 = tf.add_paragraph()
            p2.level = 1
            run2 = p2.add_run()
            run2.text = f"    {detail}"
            run2.font.size = Pt(size - 2)
            run2.font.bold = False
            run2.font.color.rgb = DARK_GRAY

    return txBox


prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank


# ─── SLIDE 1: Title ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, AMAZON_NAVY)

# Orange accent bar (left stripe)
add_rect(s1, 0, 0, Inches(0.18), SLIDE_H, AMAZON_ORANGE)

# Orange horizontal accent near bottom
add_rect(s1, Inches(0.5), Inches(5.6), Inches(7), Inches(0.06), AMAZON_ORANGE)

# Large game tile decoration (top-right)
tile_size = Inches(2.2)
for row in range(3):
    for col in range(3):
        val = 2 ** (row * 3 + col + 1)
        shades = [
            RGBColor(0x2E, 0x3A, 0x4E),
            RGBColor(0x33, 0x42, 0x59),
            RGBColor(0x1A, 0x26, 0x35),
            RGBColor(0x3D, 0x52, 0x66),
            RGBColor(0x28, 0x38, 0x4D),
            RGBColor(0x3A, 0x4E, 0x61),
            RGBColor(0x22, 0x30, 0x42),
            RGBColor(0x41, 0x56, 0x6A),
            RGBColor(0x2C, 0x3E, 0x52),
        ]
        tile = add_rect(
            s1,
            Inches(10.0) + col * Inches(0.9),
            Inches(0.3) + row * Inches(0.9),
            Inches(0.82), Inches(0.82),
            shades[row * 3 + col]
        )
        txb = slide_txb = s1.shapes.add_textbox(
            Inches(10.0) + col * Inches(0.9),
            Inches(0.3) + row * Inches(0.9),
            Inches(0.82), Inches(0.82)
        )
        tf = txb.text_frame
        tf.text = str(val)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = AMAZON_ORANGE

# "2048" big number watermark
wm = add_textbox(s1, Inches(8.5), Inches(3.5), Inches(4.5), Inches(2.5),
                 "2048", 96, bold=True, color=RGBColor(0x2E, 0x3C, 0x50),
                 align=PP_ALIGN.CENTER)

# Subtitle label
add_textbox(s1, Inches(0.6), Inches(1.5), Inches(9), Inches(0.6),
            "PROJECT OVERVIEW", 13, bold=True, color=AMAZON_ORANGE,
            align=PP_ALIGN.LEFT)

# Main title
add_textbox(s1, Inches(0.6), Inches(2.1), Inches(9), Inches(1.6),
            "2048 Game", 56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Tagline
add_textbox(s1, Inches(0.6), Inches(3.8), Inches(9), Inches(0.5),
            "Built with React · TypeScript · EPCC Workflow", 18,
            color=LIGHT_GRAY, align=PP_ALIGN.LEFT, italic=True)

# Footer
add_rect(s1, 0, Inches(6.9), SLIDE_W, Inches(0.6), RGBColor(0x1A, 0x24, 0x32))
add_textbox(s1, Inches(0.6), Inches(6.95), Inches(6), Inches(0.4),
            "Amazon Internal  |  2026", 11, color=RGBColor(0xAA, 0xAA, 0xAA),
            align=PP_ALIGN.LEFT)
add_textbox(s1, Inches(7), Inches(6.95), Inches(6), Inches(0.4),
            "1 / 3", 11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ─── SLIDE 2: Features ────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)

# White background
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, WHITE)

# Top header band
add_rect(s2, 0, 0, SLIDE_W, Inches(1.35), AMAZON_NAVY)

# Orange left stripe
add_rect(s2, 0, 0, Inches(0.18), SLIDE_H, AMAZON_ORANGE)

# Slide title
add_textbox(s2, Inches(0.5), Inches(0.28), Inches(10), Inches(0.75),
            "Key Features", 36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(s2, Inches(0.5), Inches(0.95), Inches(10), Inches(0.35),
            "What makes this implementation stand out", 14,
            color=AMAZON_ORANGE, align=PP_ALIGN.LEFT)

# Feature cards — 2×2 grid
features = [
    ("Score History",
     "Tracks high scores across sessions with persistent local storage. "
     "Players can review past performance and compare scores over time."),
    ("Leaderboard",
     "Global and local leaderboard views with real-time ranking updates. "
     "Displays top players with timestamps and score breakdowns."),
    ("Game State Persistence",
     "Auto-saves board state on every move via custom React hooks. "
     "Resume exactly where you left off — even after a page refresh."),
    ("React + TypeScript",
     "Fully type-safe component architecture with strict mode enabled. "
     "Custom hooks abstract game logic from UI for clean separation."),
]

cols = 2
card_w = Inches(5.8)
card_h = Inches(2.2)
pad_x  = Inches(0.55)
pad_y  = Inches(1.55)
gap_x  = Inches(0.35)
gap_y  = Inches(0.22)

icons = ["📊", "🏆", "💾", "⚛"]

for idx, (title, desc) in enumerate(features):
    row = idx // cols
    col = idx % cols
    x = pad_x + col * (card_w + gap_x)
    y = pad_y + row * (card_h + gap_y)

    # Card shadow (offset rectangle)
    add_rect(s2, x + Inches(0.05), y + Inches(0.05), card_w, card_h,
             RGBColor(0xDD, 0xDD, 0xDD))

    # Card background
    card = add_rect(s2, x, y, card_w, card_h, WHITE,
                    line_color=RGBColor(0xE0, 0xE0, 0xE0))

    # Orange top accent on card
    add_rect(s2, x, y, card_w, Inches(0.07), AMAZON_ORANGE)

    # Card title
    add_textbox(s2, x + Inches(0.22), y + Inches(0.15),
                card_w - Inches(0.44), Inches(0.5),
                f"{icons[idx]}  {title}", 17, bold=True,
                color=AMAZON_NAVY, align=PP_ALIGN.LEFT)

    # Card description
    add_textbox(s2, x + Inches(0.22), y + Inches(0.65),
                card_w - Inches(0.44), Inches(1.35),
                desc, 13, color=DARK_GRAY, align=PP_ALIGN.LEFT)

# Footer
add_rect(s2, 0, Inches(6.9), SLIDE_W, Inches(0.6), AMAZON_NAVY)
add_textbox(s2, Inches(0.6), Inches(6.95), Inches(6), Inches(0.4),
            "Amazon Internal  |  2026", 11, color=RGBColor(0xAA, 0xAA, 0xAA))
add_textbox(s2, Inches(7), Inches(6.95), Inches(6), Inches(0.4),
            "2 / 3", 11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ─── SLIDE 3: EPCC Summary ────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)

add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_rect(s3, 0, 0, SLIDE_W, Inches(1.35), AMAZON_NAVY)
add_rect(s3, 0, 0, Inches(0.18), SLIDE_H, AMAZON_ORANGE)

add_textbox(s3, Inches(0.5), Inches(0.28), Inches(10), Inches(0.75),
            "EPCC Workflow", 36, bold=True, color=WHITE)
add_textbox(s3, Inches(0.5), Inches(0.95), Inches(10), Inches(0.35),
            "How we built the 2048 Game end-to-end", 14, color=AMAZON_ORANGE)

# Phases
phases = [
    ("E", "Explore",  AMAZON_ORANGE,
     "Deep-dived into 2048 game rules, existing implementations, "
     "and React + TypeScript best practices before writing a single line."),
    ("P", "Plan",     RGBColor(0x00, 0x73, 0xBB),
     "Designed component hierarchy, state model, hook interfaces, "
     "and persistence strategy. Agreed on acceptance criteria up-front."),
    ("C", "Code",     RGBColor(0x1E, 0x8A, 0x44),
     "Implemented game engine, board UI, score history, leaderboard, "
     "and localStorage persistence with full TypeScript strict mode."),
    ("C", "Commit",   RGBColor(0x8A, 0x1E, 0x8A),
     "Ran typecheck, lint, and unit tests. Reviewed code quality, "
     "wrote documentation, and delivered a production-ready build."),
]

phase_w = Inches(2.85)
phase_h = Inches(3.8)
phase_y = Inches(1.7)
phase_gap = Inches(0.22)
total_w = 4 * phase_w + 3 * phase_gap
start_x = (SLIDE_W - total_w) / 2

for i, (letter, name, color, desc) in enumerate(phases):
    x = start_x + i * (phase_w + phase_gap)

    # Phase card
    add_rect(s3, x, phase_y, phase_w, phase_h, WHITE,
             line_color=RGBColor(0xE0, 0xE0, 0xE0))

    # Colored top bar
    add_rect(s3, x, phase_y, phase_w, Inches(0.08), color)

    # Circle badge
    badge = add_rect(s3, x + phase_w / 2 - Inches(0.42),
                     phase_y + Inches(0.18),
                     Inches(0.84), Inches(0.84), color)
    badge.line.fill.background()

    # Letter inside badge
    add_textbox(s3, x + phase_w / 2 - Inches(0.42),
                phase_y + Inches(0.18),
                Inches(0.84), Inches(0.84),
                letter, 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Phase name
    add_textbox(s3, x + Inches(0.15), phase_y + Inches(1.18),
                phase_w - Inches(0.3), Inches(0.5),
                name, 17, bold=True, color=AMAZON_NAVY, align=PP_ALIGN.CENTER)

    # Orange underline
    add_rect(s3, x + phase_w / 2 - Inches(0.5), phase_y + Inches(1.65),
             Inches(1.0), Inches(0.04), color)

    # Description
    add_textbox(s3, x + Inches(0.18), phase_y + Inches(1.78),
                phase_w - Inches(0.36), Inches(1.8),
                desc, 12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < 3:
        ax = x + phase_w + Inches(0.04)
        add_textbox(s3, ax, phase_y + Inches(1.55),
                    phase_gap, Inches(0.5),
                    "→", 18, bold=True,
                    color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Bottom tagline
add_rect(s3, Inches(0.5), Inches(5.7), SLIDE_W - Inches(1), Inches(0.06),
         RGBColor(0xEE, 0xEE, 0xEE))
add_textbox(s3, Inches(0.5), Inches(5.85), SLIDE_W - Inches(1), Inches(0.5),
            "Structured, intentional development — from first principles to shipped product.",
            14, italic=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Footer
add_rect(s3, 0, Inches(6.9), SLIDE_W, Inches(0.6), AMAZON_NAVY)
add_textbox(s3, Inches(0.6), Inches(6.95), Inches(6), Inches(0.4),
            "Amazon Internal  |  2026", 11, color=RGBColor(0xAA, 0xAA, 0xAA))
add_textbox(s3, Inches(7), Inches(6.95), Inches(6), Inches(0.4),
            "3 / 3", 11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


out = "/workshop/2048_Game_Project_Overview.pptx"
prs.save(out)
print(f"Saved: {out}")
